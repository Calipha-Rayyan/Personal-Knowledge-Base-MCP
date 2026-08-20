import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)
    return "\n".join(text)

def extract_text_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_text_from_md(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    elif ext == '.md':
        return extract_text_from_md(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")